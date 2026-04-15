from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTFILE = ROOT / 'gnn_ip_pair_execution_seminar.pptx'
LOGO = ROOT.parent / 'source' / '_static' / 'SevenNet_logo.png'


BG = RGBColor(248, 246, 240)
INK = RGBColor(25, 33, 41)
MUTED = RGBColor(88, 100, 117)
ACCENT = RGBColor(197, 83, 29)
ACCENT_DARK = RGBColor(125, 47, 17)
ACCENT_SOFT = RGBColor(235, 216, 203)
STEEL = RGBColor(48, 76, 97)
STEEL_SOFT = RGBColor(217, 229, 238)


SLIDES = [
    {
        'kind': 'title',
        'title': 'GNN-IP Pair Execution 제안 세미나',
        'subtitle': (
            'Equivariant GNN 기반 interatomic potential의 중복 기하 연산과 '
            'node-centric reduction 오버헤드를 줄이기 위한 실행 구조 제안'
        ),
        'footer': 'SevenNet / Proposal Seminar / 2026-03',
    },
    {
        'kind': 'bullets',
        'title': '왜 지금 이 문제를 보나',
        'bullets': [
            '현재 GNN-IP 추론은 하나의 물리적 상호작용을 양방향 directed edge로 다루기 때문에 공유 가능한 기하 정보가 반복 처리된다.',
            '문제는 단순한 tensor-product 속도 부족이 아니라 그래프를 표현하는 단위와 실행 순서가 실제 연산 구조와 어긋난다는 점에 있다.',
            '평균 이웃 수가 높은 short-range MLIP에서는 이 중복이 계산량, 메모리 트래픽, 집계 비용, 분산 통신 비용으로 동시에 확대된다.',
        ],
    },
    {
        'kind': 'pipeline',
        'title': '현재 실행 경로',
        'steps': [
            'Directed edge graph',
            'Radial / cutoff / SH',
            'weight_nn',
            'Tensor product',
            'Node sum reduction',
            'Energy / Force / Stress',
        ],
        'caption': (
            '현재 경로는 edge 중심 생성과 node 중심 집계를 분리해 두며, '
            'reverse edge에서 공유 가능한 계산도 그대로 반복한다.'
        ),
    },
    {
        'kind': 'bullets',
        'title': '코드에서 확인한 병목',
        'bullets': [
            '기하 정보 재사용은 가능하지만, 현재 구현은 reverse edge의 geometry와 filter를 생성하는 쪽이 여전히 hot path에 놓여 있다.',
            'pair metadata 생성과 topology cache 검사는 Python/C++ 모두 일반적인 lookup 구조에 의존해 control-path 비용을 남긴다.',
            'backend 선택은 아직 정적 또는 휴리스틱 수준이며, FlashTP와 pair execution을 함께 쓸 때 최적 경로가 자동으로 보장되지 않는다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '이번 세미나의 문제 정의',
        'bullets': [
            '문제 1: 동일 상호작용에서 공유 가능한 기하 정보가 directed edge 기준으로 중복 생성된다.',
            '문제 2: edge에서 생성된 message가 다시 node 기준 reduction으로 모이면서 중간 결과 저장과 메모리 이동이 커진다.',
            '문제 3: 분산 실행에서는 pair 정보와 backward 경로가 실제 필요 이상으로 일반적인 형태로 유지될 가능성이 크다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '제안 A: Exact Pair-Symmetric Execution',
        'bullets': [
            '기본 실행 단위를 directed edge가 아니라 하나의 상호작용 pair로 다시 정의한다.',
            '거리, radial basis, cutoff, pair-level SH는 한 번만 만들고 reverse 방향은 parity-aware transform으로 복원한다.',
            '최종 message 자체는 source feature가 다르므로 공유하지 않되, geometry와 filter 생성 비용은 제거 가능한 범위까지 제거한다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '제안 B: Backend Co-Design',
        'bullets': [
            'FlashTP와 같은 가속 backend에서는 pair-major layout을 가능한 오래 유지해 directed expansion 시점을 뒤로 미룬다.',
            'runtime이 backend, graph 규모, degree, device capability를 보고 baseline / geometry_only / full 중 최적 경로를 고르게 한다.',
            '핵심은 공통 옵션을 유지하되, backend마다 같은 구현을 억지로 쓰지 않고 가장 빠른 하위 경로를 선택하도록 만드는 것이다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '제안 C: Topology-Epoch Caching',
        'bullets': [
            'MD에서는 좌표는 자주 바뀌어도 neighbor topology는 일정 구간 유지될 수 있다.',
            '이 구간에서는 pair plan, reverse map, parity metadata, reduction schedule, distributed comm metadata를 재사용한다.',
            '핵심은 단순 캐시가 아니라, 어떤 이벤트에서 무효화되는지까지 포함한 topology-epoch 실행 모델을 정의하는 것이다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '제안 D: Distributed Pair Schedule',
        'bullets': [
            '분산 경로에서는 ghost 정보를 무조건 많이 들고 가는 대신, pair-aware schedule과 partial reduction 관점으로 통신을 다시 본다.',
            '특히 scalar energy readout에 실제로 기여하는 경로만 남기는 scalar-only backward pruning 가능성을 본다.',
            '이 방향은 단순 partitioning이 아니라, exactness를 유지하면서 통신되는 정보의 구조를 바꾸는 제안이다.',
        ],
    },
    {
        'kind': 'twocol',
        'title': '기대효과와 리스크',
        'left_title': '기대효과',
        'left_items': [
            '중복 기하 연산 감소',
            '메모리 트래픽 완화',
            'backend별 최적 실행 경로 확보',
            '분산 경로의 통신 구조 재설계 가능성',
        ],
        'right_title': '리스크',
        'right_items': [
            'pair metadata 생성 자체가 새로운 병목이 될 수 있음',
            'FlashTP 등 가속 backend에서 조기 expansion이 이득을 상쇄할 수 있음',
            'distributed exactness 증명과 backward pruning 조건 정리가 필요함',
            'topology cache invalidation을 잘못 정의하면 correctness가 깨질 수 있음',
        ],
    },
    {
        'kind': 'bullets',
        'title': '검증 계획',
        'bullets': [
            '정확성: energy, atomic energy, force, stress parity와 long-run MD 통계량을 함께 확인한다.',
            '성능: baseline / geometry_only / full / backend auto policy를 단일 GPU와 분산 경로에서 모두 비교한다.',
            'ablation: pair execution, fused reduction, topology cache, distributed pruning을 독립적으로 떼어 효과를 분리한다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '세미나에서 가져갈 핵심 메시지',
        'bullets': [
            '이 연구는 새 MLIP 모델 제안이 아니라, 기존 equivariant MLIP의 실행 구조를 exact하게 재구성하는 연구다.',
            '핵심 novelty 후보는 pair execution 그 자체보다도 backend co-design, topology-epoch runtime, distributed exact pruning에 있다.',
            '즉, “어떤 연산자를 빠르게 만들 것인가”보다 “어떤 단위로 표현하고 어떤 시점에 계산할 것인가”가 중심 질문이다.',
        ],
    },
    {
        'kind': 'bullets',
        'title': '마무리',
        'bullets': [
            '단기 목표: pair plan hot path 제거와 backend autotuning 기준 정립',
            '중기 목표: FlashTP pair-major adapter와 topology-epoch caching 정식화',
            '장기 목표: distributed exact pruning까지 포함한 시스템 논문 스토리 완성',
        ],
    },
]


SCRIPT = [
    {
        'title': 'Slide 1. 제목',
        'script': (
            '오늘 세미나는 구현 결과 보고가 아니라, GNN-IP 추론 실행 구조를 어떻게 다시 정의할지에 대한 제안 세미나입니다. '
            '핵심 문제는 directed edge 기준 실행이 실제 물리적 상호작용 구조와 맞지 않아 중복 계산과 집계 오버헤드를 만든다는 점입니다.'
        ),
    },
    {
        'title': 'Slide 2. 왜 지금 이 문제를 보나',
        'script': (
            'short-range MLIP는 그래프 이론 의미의 dense graph는 아니지만, local degree가 높고 기하적 규칙성이 강합니다. '
            '이런 구조에서는 상호작용을 두 방향 간선으로 그대로 실행하는 것이 자연스러운 구현일 수는 있어도 효율적인 구현이라고 보기 어렵습니다.'
        ),
    },
    {
        'title': 'Slide 3. 현재 실행 경로',
        'script': (
            '현재 실행은 edge 생성, 기하 feature 생성, weight_nn, tensor product, node reduction, 그리고 energy-force-stress 계산으로 이어집니다. '
            '즉 edge 중심 생성과 node 중심 reduction이 분리되어 있고, 이 구조가 중복과 메모리 이동을 함께 키웁니다.'
        ),
    },
    {
        'title': 'Slide 4. 코드에서 확인한 병목',
        'script': (
            '실제 코드에서도 pair metadata 생성, topology cache 검증, backend 선택이 아직 충분히 공격적으로 최적화되어 있지 않습니다. '
            '그래서 단순히 연산 수를 줄이는 것만으로는 부족하고, control path와 backend path를 함께 재설계해야 합니다.'
        ),
    },
    {
        'title': 'Slide 5. 문제 정의',
        'script': (
            '정리하면 문제는 세 가지입니다. 공유 가능한 기하 정보의 중복 생성, message 이후 node-centric reduction에 따른 메모리 비용, '
            '그리고 분산 경로에서 과도하게 일반적인 통신 구조입니다.'
        ),
    },
    {
        'title': 'Slide 6. Exact Pair-Symmetric Execution',
        'script': (
            '첫 번째 축은 pair를 기본 실행 단위로 보는 것입니다. '
            'geometry와 filter는 pair 기준으로 만들고, reverse 방향은 parity transform으로 복원합니다. '
            '다만 source feature가 다르므로 message 자체는 공유하지 않는다는 점을 명확히 구분해야 합니다.'
        ),
    },
    {
        'title': 'Slide 7. Backend Co-Design',
        'script': (
            '두 번째 축은 backend co-design입니다. 공통 옵션을 둔다고 해서 backend 내부까지 똑같이 실행하면 안 됩니다. '
            '특히 FlashTP에서는 pair-major layout을 오래 유지하고, 실제 측정 기반으로 policy를 고르는 autotuning이 필요합니다.'
        ),
    },
    {
        'title': 'Slide 8. Topology-Epoch Caching',
        'script': (
            '세 번째 축은 topology-epoch caching입니다. '
            'MD에서는 neighbor topology가 유지되는 구간이 있기 때문에, pair plan과 reduction schedule을 이 구간 단위로 재사용하는 것이 가능합니다. '
            '핵심은 무엇을 cache하느냐보다 어떤 조건에서 invalidate하느냐입니다.'
        ),
    },
    {
        'title': 'Slide 9. Distributed Pair Schedule',
        'script': (
            '네 번째 축은 distributed exactness입니다. '
            'ghost를 많이 들고 가는 구조 대신, pair-aware partial reduction과 scalar-only backward pruning이 가능하다면 통신 구조 자체를 바꿀 수 있습니다.'
        ),
    },
    {
        'title': 'Slide 10. 기대효과와 리스크',
        'script': (
            '이 연구의 장점은 계산량, 메모리, backend utilization, distributed scaling을 함께 다룰 수 있다는 점입니다. '
            '반면 pair metadata 생성이 새 병목이 되거나, backend별로 pair execution 이득이 상쇄될 수 있다는 위험도 분명히 있습니다.'
        ),
    },
    {
        'title': 'Slide 11. 검증 계획',
        'script': (
            '검증은 정확성, 성능, ablation 세 축으로 가져갑니다. '
            '특히 baseline과 수치적으로 같은지, 그리고 어떤 제안이 실제 속도 향상에 기여하는지를 분리해서 봐야 합니다.'
        ),
    },
    {
        'title': 'Slide 12. 핵심 메시지',
        'script': (
            '결국 이 연구는 새 모델 제안이 아니라 exact runtime 재구성 연구입니다. '
            '핵심 질문은 빠른 커널 하나가 아니라, 어떤 단위로 표현하고 언제 계산하며 무엇을 재사용할 것인가입니다.'
        ),
    },
    {
        'title': 'Slide 13. 마무리',
        'script': (
            '단기적으로는 pair plan hot path와 backend autotuning을, 중기적으로는 FlashTP adapter와 topology-epoch runtime을, '
            '장기적으로는 distributed exact pruning까지 포함한 시스템 논문 스토리를 목표로 가져가면 됩니다.'
        ),
    },
]


def _add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_footer(slide, text: str, slide_no: int):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.0), Inches(12.3), Inches(0.25))
    p = box.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = f'{text}    {slide_no}'
    p.alignment = PP_ALIGN.RIGHT
    run.font.size = Pt(10)
    run.font.color.rgb = MUTED


def _style_title(shape):
    p = shape.text_frame.paragraphs[0]
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK


def _add_bullet_slide(prs, item, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.0), Inches(0.6))
    title.text = item['title']
    _style_title(title)

    accent = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.7),
        Inches(1.15),
        Inches(1.1),
        Inches(0.08),
    )
    accent.fill.solid()
    accent.fill.fore_color.rgb = ACCENT
    accent.line.fill.background()

    body = slide.shapes.add_textbox(Inches(0.9), Inches(1.55), Inches(11.0), Inches(4.9))
    tf = body.text_frame
    tf.word_wrap = True
    for idx, bullet in enumerate(item['bullets']):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.space_after = Pt(14)
        p.font.size = Pt(22)
        p.font.color.rgb = INK

    _add_footer(slide, 'Proposal Seminar', slide_no)
    return slide


def _add_title_slide(prs, item, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.0),
        Inches(0.0),
        Inches(13.333),
        Inches(1.0),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = STEEL
    band.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.85), Inches(1.45), Inches(10.8), Inches(1.2))
    p = title.text_frame.paragraphs[0]
    run = p.add_run()
    run.text = item['title']
    run.font.size = Pt(30)
    run.font.bold = True
    run.font.color.rgb = INK

    subtitle = slide.shapes.add_textbox(Inches(0.9), Inches(2.65), Inches(10.7), Inches(2.0))
    p2 = subtitle.text_frame.paragraphs[0]
    p2.text = item['subtitle']
    p2.font.size = Pt(20)
    p2.font.color.rgb = MUTED

    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(10.9), Inches(1.45), width=Inches(1.7))

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(5.25),
        Inches(4.9),
        Inches(0.9),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_SOFT
    card.line.fill.background()
    tf = card.text_frame
    p3 = tf.paragraphs[0]
    p3.text = '핵심 키워드: Pair Execution / Backend Co-Design / Topology Epoch'
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_DARK

    _add_footer(slide, item['footer'], slide_no)
    return slide


def _add_pipeline_slide(prs, item, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.0), Inches(0.6))
    title.text = item['title']
    _style_title(title)

    left = 0.75
    top = 2.0
    width = 1.8
    height = 0.8
    gap = 0.25
    for idx, step in enumerate(item['steps']):
        x = Inches(left + idx * (width + gap))
        shape = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            x,
            Inches(top),
            Inches(width),
            Inches(height),
        )
        shape.fill.solid()
        shape.fill.fore_color.rgb = ACCENT_SOFT if idx % 2 == 0 else STEEL_SOFT
        shape.line.color.rgb = ACCENT if idx % 2 == 0 else STEEL
        tf = shape.text_frame
        p = tf.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INK
        if idx < len(item['steps']) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(left + idx * (width + gap) + width + 0.02),
                Inches(top + 0.18),
                Inches(0.22),
                Inches(0.42),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    caption = slide.shapes.add_textbox(Inches(0.95), Inches(4.1), Inches(11.0), Inches(1.6))
    p = caption.text_frame.paragraphs[0]
    p.text = item['caption']
    p.font.size = Pt(20)
    p.font.color.rgb = MUTED
    _add_footer(slide, 'Proposal Seminar', slide_no)
    return slide


def _add_twocol_slide(prs, item, slide_no):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    title = slide.shapes.add_textbox(Inches(0.7), Inches(0.45), Inches(11.0), Inches(0.6))
    title.text = item['title']
    _style_title(title)

    for x, box_title, entries, fill_rgb in (
        (0.8, item['left_title'], item['left_items'], STEEL_SOFT),
        (6.7, item['right_title'], item['right_items'], ACCENT_SOFT),
    ):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.5),
            Inches(5.1),
            Inches(4.7),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill_rgb
        box.line.fill.background()
        tx = slide.shapes.add_textbox(Inches(x + 0.2), Inches(1.75), Inches(4.6), Inches(0.4))
        p = tx.text_frame.paragraphs[0]
        p.text = box_title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = INK
        body = slide.shapes.add_textbox(Inches(x + 0.28), Inches(2.35), Inches(4.4), Inches(3.5))
        tf = body.text_frame
        for idx, bullet in enumerate(entries):
            p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
            p.text = bullet
            p.font.size = Pt(18)
            p.font.color.rgb = INK
            p.space_after = Pt(12)

    _add_footer(slide, 'Proposal Seminar', slide_no)
    return slide


def build_presentation():
    prs = Presentation()
    prs.slide_width = Inches(13.333)
    prs.slide_height = Inches(7.5)

    for idx, item in enumerate(SLIDES, start=1):
        kind = item['kind']
        if kind == 'title':
            _add_title_slide(prs, item, idx)
        elif kind == 'pipeline':
            _add_pipeline_slide(prs, item, idx)
        elif kind == 'twocol':
            _add_twocol_slide(prs, item, idx)
        else:
            _add_bullet_slide(prs, item, idx)

    prs.save(str(OUTFILE))


if __name__ == '__main__':
    build_presentation()
