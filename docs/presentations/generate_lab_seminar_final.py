from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_AUTO_SIZE, PP_ALIGN
from pptx.util import Inches, Pt


ROOT = Path(__file__).resolve().parent
OUTFILE = ROOT / 'Seminar_김민창_GNN-IP_최종세미나.pptx'
SCRIPT_OUT = ROOT / 'Seminar_김민창_GNN-IP_최종세미나_대본.md'
LOGO = ROOT.parent / 'source' / '_static' / 'SevenNet_logo.png'

ASSET_ROOT = ROOT / 'assets' / 'quadrant_pack'
PLOT_MAP = ASSET_ROOT / 'plots' / 'quadrant_dataset_map.png'
PLOT_LAT = ASSET_ROOT / 'plots' / 'quadrant_latency_comparison.png'
PLOT_SPEED = ASSET_ROOT / 'plots' / 'quadrant_speedup_comparison.png'
PLOT_STAGE = ASSET_ROOT / 'plots' / 'extreme_stage_breakdown.png'
PLOT_DOWNLOAD = ASSET_ROOT / 'plots' / 'large_dataset_download_status.png'
DIAGRAM_QUADRANT = ASSET_ROOT / 'diagrams' / 'quadrant_mechanism_diagram.png'


BG = RGBColor(248, 246, 240)
INK = RGBColor(28, 34, 41)
MUTED = RGBColor(90, 99, 112)
ACCENT = RGBColor(191, 77, 32)
ACCENT_DARK = RGBColor(130, 49, 22)
ACCENT_SOFT = RGBColor(237, 221, 209)
STEEL = RGBColor(45, 72, 97)
STEEL_SOFT = RGBColor(219, 230, 238)
GREEN_SOFT = RGBColor(218, 240, 229)
GOLD_SOFT = RGBColor(247, 236, 206)


SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)


NOTES = [
    {
        'title': 'Slide 1. 제목',
        'time': '1.0 min',
        'script': (
            '오늘 발표는 GNN-IP 추론 최적화, 그중에서도 pair execution 관점에서 지금까지 무엇을 이해했고 '
            '무엇을 구현했고, 어떤 실험 결과가 나왔는지를 정리하는 랩 세미나입니다. '
            '핵심은 새 모델 제안이 아니라, equivariant interatomic potential의 실행 구조를 다시 보는 것입니다.'
        ),
    },
    {
        'title': 'Slide 2. 발표 구성',
        'time': '1.0 min',
        'script': (
            '발표는 네 부분으로 진행하겠습니다. 먼저 GNN-IP 전체를 짧게 설명하고, '
            '그 위에서 재사용 가능한 연산값을 어떻게 분류했는지 보겠습니다. '
            '그 다음 SevenNet 기준 구현과 관련 최적화 연구를 연결하고, 마지막으로 실험 결과와 discussion으로 마무리하겠습니다.'
        ),
    },
    {
        'title': 'Slide 3. GNN-IP란 무엇인가',
        'time': '2.0 min',
        'script': (
            '작업자 외 학생들을 위해 가장 짧게 정리하면, GNN-IP는 원자 구조를 그래프로 보고 '
            '원자별 local environment로부터 total energy를 예측한 뒤, force와 stress를 gradient로 얻는 모델입니다. '
            '핵심은 원자 위치가 회전해도 물리량이 맞게 변해야 하므로 equivariance가 중요하다는 점입니다.'
        ),
    },
    {
        'title': 'Slide 4. Equivariant 추론 파이프라인',
        'time': '2.0 min',
        'script': (
            '실제 추론은 atom과 neighbor list에서 시작해 edge geometry를 만들고, radial basis, cutoff, '
            'spherical harmonics로 edge feature를 만든 뒤, 여러 개의 interaction block에서 weight_nn, tensor product, '
            'aggregation을 반복합니다. 마지막에 energy를 읽고 force와 stress를 계산합니다.'
        ),
    },
    {
        'title': 'Slide 5. 재사용 연산값의 분류',
        'time': '2.0 min',
        'script': (
            '지난 주 계획서의 첫 항목은 재사용 가능한 값을 분류하는 것이었습니다. '
            '여기서 중요한 구분은 pair-shared geometry, parity로 복원 가능한 값, 방향별로 다시 계산해야 하는 message, '
            '그리고 node/global reduction 결과입니다. 이 분류가 있어야 어디까지 exact하게 줄일 수 있는지가 보입니다.'
        ),
    },
    {
        'title': 'Slide 6. SevenNet 현재 구현과 목표 구조',
        'time': '2.5 min',
        'script': (
            '현재 SevenNet 구현은 pair metadata를 만들고 geometry와 weight 일부를 재사용하지만, '
            '최종 TP와 aggregation은 여전히 directed edge 기준입니다. '
            '제가 최종적으로 노리는 구조는 pair-major layout을 오래 유지하고, 가능하면 pair-major TP까지 가는 구조입니다.'
        ),
    },
    {
        'title': 'Slide 7. 관련 연구',
        'time': '2.0 min',
        'script': (
            '관련 연구는 세 범주만 보면 됩니다. 첫째는 FlashTP나 cuEquivariance처럼 tensor product 자체를 빠르게 만드는 backend 연구, '
            '둘째는 graph/runtime 측면의 cache와 batching, 셋째는 distributed inference/runtime 연구입니다. '
            '이번 발표는 이 세 범주 중 pair execution과 직접 연결되는 부분만 가져왔습니다.'
        ),
    },
    {
        'title': 'Slide 8. 실험 질문과 측정 항목',
        'time': '1.5 min',
        'script': (
            '실험 질문은 간단합니다. 어떤 그래프에서 pair execution이 이득인지, 왜 그런지, '
            '그리고 현재 구현이 줄이는 연산이 실제 전체 시간에서 얼마나 큰 비중을 차지하는지입니다. '
            '그래서 latency, speedup, stage profiling, GPU utilization, 데이터셋 size-density를 같이 봤습니다.'
        ),
    },
    {
        'title': 'Slide 9. 데이터셋 맵',
        'time': '2.0 min',
        'script': (
            '이 그림은 대표 샘플 기준으로 edge 수와 평균 이웃 수를 그린 것입니다. '
            '발표용으로는 size와 density를 동시에 보게 하는 게 중요해서, large 기준은 대표 샘플 `num_edges >= 3000`, '
            'dense 기준은 `avg_neighbors >= 40`으로 잡았습니다.'
        ),
    },
    {
        'title': 'Slide 10. 네 구역 대표 비교',
        'time': '2.5 min',
        'script': (
            '네 구역 대표는 small sparse로 SPICE 2023, small dense로 phononDB, large sparse로 OMol25 validation, '
            'large dense로 MPtrj validation을 잡았습니다. 이 비교가 중요한 이유는 pair execution의 성패가 '
            '크기만이 아니라 density와 edge load의 조합으로 결정된다는 점을 한 장에서 보여주기 때문입니다.'
        ),
    },
    {
        'title': 'Slide 11. Quadrant 결과 해석',
        'time': '2.0 min',
        'script': (
            '결과를 보면 small sparse는 느리고, large dense는 가장 크게 빨라집니다. '
            '재미있는 점은 small dense인 phononDB처럼 atom 수는 작아도 degree가 매우 높으면 이득이 난다는 것입니다. '
            '즉 본질 변수는 atom 수 단독이 아니라 결국 edge load입니다.'
        ),
    },
    {
        'title': 'Slide 12. Baseline 상세 프로파일',
        'time': '2.5 min',
        'script': (
            'baseline intrusive profile을 보면 small sparse에서는 message TP와 force_output이 지배적이고, '
            'large dense에서는 spherical harmonics까지 크게 올라옵니다. '
            '반면 weight_nn과 gather, aggregation은 baseline에서도 상대적으로 작습니다. '
            '그래서 현재 pair execution이 줄이는 부분만으로는 모든 그래프에서 큰 이득이 나기 어렵습니다.'
        ),
    },
    {
        'title': 'Slide 13. Lab benchmark 결과',
        'time': '2.5 min',
        'script': (
            '랩 단일 GPU 환경에서 본 실제 결과를 요약하면, large dense 쪽에서는 speedup이 확실히 보이고, '
            'small sparse에서는 손해가 납니다. 또 size-density map과 stage breakdown을 같이 보면, '
            '왜 어떤 경우는 geometry reuse가 충분히 먹히고 어떤 경우는 TP가 그대로 남아 발목을 잡는지가 같이 설명됩니다.'
        ),
    },
    {
        'title': 'Slide 14. Discussion',
        'time': '2.5 min',
        'script': (
            '현재 구현은 pair-major TP kernel이 아니라 geometry/SH/weight reuse 중심의 구현입니다. '
            '그래서 speedup ceiling이 있고, dense-large에서만 확실히 듣습니다. '
            '반대로 말하면 next step은 명확합니다. pair-major TP, backend별 policy, topology epoch cache, distributed schedule입니다.'
        ),
    },
    {
        'title': 'Slide 15. 지금까지 구현한 것과 남은 것',
        'time': '2.0 min',
        'script': (
            '이미 구현한 것은 pair metadata, geometry reuse, weight reuse, baseline/pair profiling, dataset별 실험 인프라입니다. '
            '남은 것은 pair-major TP, FlashTP와의 진짜 co-design, distributed path, 그리고 더 큰 public dataset 완전 확보입니다.'
        ),
    },
    {
        'title': 'Slide 16. 결론',
        'time': '2.0 min',
        'script': (
            '정리하면 이번 단계의 가장 중요한 결론은 세 가지입니다. '
            '첫째, pair execution을 논하려면 reusable value를 정확히 분류해야 합니다. '
            '둘째, 현재 구현은 일부 edge-side work만 줄이기 때문에 dense-large에서만 강합니다. '
            '셋째, 다음 연구 포인트는 pair-major execution과 backend/runtime co-design입니다.'
        ),
    },
    {
        'title': 'Appendix. 데이터셋 확보 상태',
        'time': 'backup',
        'script': (
            '이 슬라이드는 Q&A용입니다. repo raw 기준으로는 OC20/OC22, OMol25, sAlex train 같은 큰 셋이 아직 pending입니다. '
            '따라서 앞으로는 실험 체계를 더 밀기 전에 데이터셋 확보 상태를 정리하는 작업도 병행해야 합니다.'
        ),
    },
]


def _add_background(slide):
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def _add_footer(slide, text: str, slide_no: int):
    box = slide.shapes.add_textbox(Inches(0.45), Inches(7.03), Inches(12.3), Inches(0.24))
    p = box.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    run = p.add_run()
    run.text = f'{text}    {slide_no}'
    run.font.size = Pt(9.5)
    run.font.color.rgb = MUTED


def _title_box(slide, title: str, subtitle: str | None = None):
    tb = slide.shapes.add_textbox(Inches(0.72), Inches(0.42), Inches(11.1), Inches(0.56))
    p = tb.text_frame.paragraphs[0]
    p.text = title
    p.font.size = Pt(27)
    p.font.bold = True
    p.font.color.rgb = INK

    line = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0.72),
        Inches(1.08),
        Inches(1.18),
        Inches(0.07),
    )
    line.fill.solid()
    line.fill.fore_color.rgb = ACCENT
    line.line.fill.background()

    if subtitle:
        sb = slide.shapes.add_textbox(Inches(2.05), Inches(0.54), Inches(10.0), Inches(0.30))
        p2 = sb.text_frame.paragraphs[0]
        p2.text = subtitle
        p2.font.size = Pt(12)
        p2.font.color.rgb = MUTED


def _set_text_frame(tf, lines, font_size=22, color=INK, bullet=True):
    tf.word_wrap = True
    tf.auto_size = MSO_AUTO_SIZE.TEXT_TO_FIT_SHAPE
    for idx, line in enumerate(lines):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = line
        p.font.size = Pt(font_size)
        p.font.color.rgb = color
        p.space_after = Pt(10)
        if bullet:
            p.level = 0


def _add_logo(slide):
    if LOGO.exists():
        slide.shapes.add_picture(str(LOGO), Inches(11.2), Inches(0.36), width=Inches(1.35))


def _add_title_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)

    band = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.RECTANGLE,
        Inches(0),
        Inches(0),
        Inches(13.333),
        Inches(0.92),
    )
    band.fill.solid()
    band.fill.fore_color.rgb = STEEL
    band.line.fill.background()

    title = slide.shapes.add_textbox(Inches(0.82), Inches(1.38), Inches(8.8), Inches(1.25))
    p = title.text_frame.paragraphs[0]
    p.text = 'GNN-IP Pair Execution\n랩 세미나 최종 발표'
    p.font.size = Pt(28)
    p.font.bold = True
    p.font.color.rgb = INK

    sub = slide.shapes.add_textbox(Inches(0.88), Inches(2.86), Inches(8.8), Inches(1.1))
    p2 = sub.text_frame.paragraphs[0]
    p2.text = (
        '재사용 연산값 분류, equivariant 추론 파이프라인 이해, 관련 추론 최적화 연구, '
        '그리고 SevenNet 기준 실험 결과와 discussion'
    )
    p2.font.size = Pt(18)
    p2.font.color.rgb = MUTED

    card = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(5.28),
        Inches(5.35),
        Inches(0.78),
    )
    card.fill.solid()
    card.fill.fore_color.rgb = ACCENT_SOFT
    card.line.fill.background()
    tf = card.text_frame
    p3 = tf.paragraphs[0]
    p3.text = '김민창 / DenseMLIP / 30 min seminar'
    p3.font.size = Pt(16)
    p3.font.color.rgb = ACCENT_DARK

    _add_logo(slide)
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_agenda_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, '오늘 발표 구성', '30분')
    _add_logo(slide)

    items = [
        ('01', 'GNN-IP와 equivariant 추론 파이프라인', STEEL_SOFT),
        ('02', '재사용 가능한 연산값의 분류와 구현 포인트', ACCENT_SOFT),
        ('03', '연관 추론 최적화 연구와 현재 구현 위치', GREEN_SOFT),
        ('04', '실험 결과, profiling, discussion', GOLD_SOFT),
    ]
    y = 1.65
    for idx, (num, text, fill_rgb) in enumerate(items):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.92),
            Inches(y + idx * 1.18),
            Inches(11.1),
            Inches(0.84),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill_rgb
        box.line.fill.background()
        nbox = slide.shapes.add_textbox(Inches(1.18), Inches(y + 0.11 + idx * 1.18), Inches(0.65), Inches(0.35))
        p = nbox.text_frame.paragraphs[0]
        p.text = num
        p.font.size = Pt(18)
        p.font.bold = True
        p.font.color.rgb = ACCENT_DARK
        tbox = slide.shapes.add_textbox(Inches(2.0), Inches(y + 0.08 + idx * 1.18), Inches(9.2), Inches(0.42))
        p2 = tbox.text_frame.paragraphs[0]
        p2.text = text
        p2.font.size = Pt(22)
        p2.font.color.rgb = INK

    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_bullet_slide(prs, slide_no: int, title: str, bullets: list[str], note: str | None = None):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, title)
    _add_logo(slide)
    body = slide.shapes.add_textbox(Inches(0.95), Inches(1.5), Inches(11.05), Inches(4.95))
    _set_text_frame(body.text_frame, bullets, font_size=21, color=INK, bullet=True)
    if note:
        nb = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(0.95),
            Inches(6.22),
            Inches(11.05),
            Inches(0.58),
        )
        nb.fill.solid()
        nb.fill.fore_color.rgb = STEEL_SOFT
        nb.line.fill.background()
        nt = nb.text_frame.paragraphs[0]
        nt.text = note
        nt.font.size = Pt(13)
        nt.font.color.rgb = STEEL
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_pipeline_overview(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, 'Equivariant GNN-IP 추론 파이프라인')
    _add_logo(slide)

    steps = [
        'Atoms / Cell / Neighbor list',
        'Edge geometry\n(r, |r|, SH, basis)',
        'Equivariant\ninteraction blocks',
        'Atomic energy\nreadout',
        'Force / Stress\nvia gradient',
    ]
    x0 = 0.82
    y = 2.35
    width = 2.12
    gap = 0.30
    fills = [STEEL_SOFT, ACCENT_SOFT, GREEN_SOFT, GOLD_SOFT, STEEL_SOFT]
    for i, step in enumerate(steps):
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x0 + i * (width + gap)),
            Inches(y),
            Inches(width),
            Inches(1.15),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fills[i]
        box.line.color.rgb = STEEL
        p = box.text_frame.paragraphs[0]
        p.text = step
        p.alignment = PP_ALIGN.CENTER
        p.font.size = Pt(16)
        p.font.bold = True
        p.font.color.rgb = INK
        if i < len(steps) - 1:
            arrow = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.CHEVRON,
                Inches(x0 + i * (width + gap) + width + 0.04),
                Inches(y + 0.28),
                Inches(0.22),
                Inches(0.46),
            )
            arrow.fill.solid()
            arrow.fill.fore_color.rgb = MUTED
            arrow.line.fill.background()

    cap = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(4.55),
        Inches(10.9),
        Inches(1.3),
    )
    cap.fill.solid()
    cap.fill.fore_color.rgb = ACCENT_SOFT
    cap.line.fill.background()
    p2 = cap.text_frame.paragraphs[0]
    p2.text = (
        '핵심 포인트: energy는 node-level scalar readout이지만, 그 전에 대부분의 비용은 edge geometry 생성과 '
        'equivariant interaction block에 들어간다.'
    )
    p2.font.size = Pt(18)
    p2.font.color.rgb = INK

    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_classification_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, '재사용 가능한 연산값의 분류')
    _add_logo(slide)

    cards = [
        ('Pair-shared geometry', ['distance |r|', 'radial basis', 'cutoff'], STEEL_SOFT, 0.92, 1.6),
        ('Parity-recoverable', ['reverse SH via (-1)^l', 'canonical pair direction'], ACCENT_SOFT, 6.65, 1.6),
        ('Direction-specific', ['source feature x_i / x_j', 'final message TP output'], GREEN_SOFT, 0.92, 4.0),
        ('Node / global outputs', ['aggregation', 'atomic energy', 'force / stress'], GOLD_SOFT, 6.65, 4.0),
    ]
    for title, items, fill_rgb, x, y in cards:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(y),
            Inches(5.35),
            Inches(1.85),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill_rgb
        box.line.color.rgb = STEEL
        tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(y + 0.12), Inches(5.0), Inches(0.3))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = INK
        bb = slide.shapes.add_textbox(Inches(x + 0.24), Inches(y + 0.55), Inches(4.8), Inches(1.0))
        _set_text_frame(bb.text_frame, items, font_size=17, color=INK, bullet=True)

    note = slide.shapes.add_shape(
        MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
        Inches(1.0),
        Inches(6.28),
        Inches(10.9),
        Inches(0.44),
    )
    note.fill.solid()
    note.fill.fore_color.rgb = ACCENT_SOFT
    note.line.fill.background()
    p2 = note.text_frame.paragraphs[0]
    p2.text = '이 분류가 있어야 “무엇을 exact하게 줄일 수 있고, 무엇은 줄일 수 없는가”를 구분할 수 있다.'
    p2.font.size = Pt(15)
    p2.font.color.rgb = ACCENT_DARK

    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_current_vs_target_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, 'SevenNet 현재 구현과 목표 구조')
    _add_logo(slide)

    cols = [
        (
            0.86,
            '현재 구현',
            [
                'pair metadata',
                'geometry / SH 일부 재사용',
                'weight_nn 일부 재사용',
                'directed TP',
                'node aggregation',
            ],
            STEEL_SOFT,
            '현재는 pair-aware reuse이지, pair-major TP는 아니다.',
        ),
        (
            6.72,
            '목표 구조',
            [
                'pair metadata',
                'pair-major geometry 유지',
                'pair-major TP or delayed expansion',
                'backend-aware reduction',
                'runtime policy selection',
            ],
            ACCENT_SOFT,
            '핵심은 pair layout을 끝까지 유지하는 것이다.',
        ),
    ]
    for x, title, steps, fill_rgb, caption in cols:
        card = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.55),
            Inches(5.12),
            Inches(4.75),
        )
        card.fill.solid()
        card.fill.fore_color.rgb = fill_rgb
        card.line.color.rgb = STEEL

        tt = slide.shapes.add_textbox(Inches(x + 0.18), Inches(1.72), Inches(4.6), Inches(0.30))
        p = tt.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(22)
        p.font.bold = True
        p.font.color.rgb = INK

        y = 2.35
        for idx, step in enumerate(steps):
            b = slide.shapes.add_shape(
                MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
                Inches(x + 0.28),
                Inches(y + idx * 0.72),
                Inches(4.46),
                Inches(0.48),
            )
            b.fill.solid()
            b.fill.fore_color.rgb = BG
            b.line.color.rgb = MUTED
            t = b.text_frame.paragraphs[0]
            t.text = step
            t.alignment = PP_ALIGN.CENTER
            t.font.size = Pt(14)
            t.font.bold = True
            t.font.color.rgb = INK

        cap = slide.shapes.add_textbox(Inches(x + 0.22), Inches(5.78), Inches(4.6), Inches(0.36))
        p2 = cap.text_frame.paragraphs[0]
        p2.text = caption
        p2.font.size = Pt(13)
        p2.font.color.rgb = MUTED

    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_related_work_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, '연관 추론 최적화 연구: 이 발표에 필요한 부분만')
    _add_logo(slide)

    blocks = [
        ('Backend fusion', ['FlashTP', 'cuEquivariance / OEQ', 'TP kernel optimization'], STEEL_SOFT, 0.92),
        ('Runtime / graph', ['neighbor cache', 'graph layout', 'batching / memory traffic'], ACCENT_SOFT, 4.4),
        ('System / distributed', ['ghost communication', 'schedule', 'runtime policy'], GREEN_SOFT, 7.88),
    ]
    for title, items, fill_rgb, x in blocks:
        box = slide.shapes.add_shape(
            MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE,
            Inches(x),
            Inches(1.9),
            Inches(3.2),
            Inches(3.45),
        )
        box.fill.solid()
        box.fill.fore_color.rgb = fill_rgb
        box.line.color.rgb = STEEL
        tb = slide.shapes.add_textbox(Inches(x + 0.18), Inches(2.08), Inches(2.8), Inches(0.4))
        p = tb.text_frame.paragraphs[0]
        p.text = title
        p.font.size = Pt(20)
        p.font.bold = True
        p.font.color.rgb = INK
        body = slide.shapes.add_textbox(Inches(x + 0.20), Inches(2.68), Inches(2.75), Inches(2.2))
        _set_text_frame(body.text_frame, items, font_size=17, color=INK, bullet=True)

    note = slide.shapes.add_textbox(Inches(1.0), Inches(6.0), Inches(11.0), Inches(0.5))
    p2 = note.text_frame.paragraphs[0]
    p2.text = '이번 발표는 커널 최적화 자체보다, pair-aware runtime과 backend co-design이 어디에 들어가는지를 설명하는 데 초점을 둔다.'
    p2.font.size = Pt(16)
    p2.font.color.rgb = MUTED
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_image_slide(prs, slide_no: int, title: str, image: Path, bullets: list[str] | None = None, image_right: bool = False):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, title)
    _add_logo(slide)
    if image_right:
        img_x, txt_x = 6.5, 0.92
    else:
        img_x, txt_x = 0.92, 8.95 if bullets else 0.92

    if image.exists():
        slide.shapes.add_picture(str(image), Inches(img_x), Inches(1.45), width=Inches(5.6 if bullets else 11.0), height=Inches(4.95 if bullets else 5.6))
    if bullets:
        body = slide.shapes.add_textbox(Inches(txt_x), Inches(1.75), Inches(4.0), Inches(4.6))
        _set_text_frame(body.text_frame, bullets, font_size=18, color=INK, bullet=True)
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_two_image_slide(prs, slide_no: int, title: str, left_img: Path, right_img: Path, left_title: str, right_title: str):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, title)
    _add_logo(slide)

    lt = slide.shapes.add_textbox(Inches(0.98), Inches(1.28), Inches(5.2), Inches(0.28))
    p1 = lt.text_frame.paragraphs[0]
    p1.text = left_title
    p1.font.size = Pt(16)
    p1.font.bold = True
    p1.font.color.rgb = INK
    rt = slide.shapes.add_textbox(Inches(6.82), Inches(1.28), Inches(5.2), Inches(0.28))
    p2 = rt.text_frame.paragraphs[0]
    p2.text = right_title
    p2.font.size = Pt(16)
    p2.font.bold = True
    p2.font.color.rgb = INK

    if left_img.exists():
        slide.shapes.add_picture(str(left_img), Inches(0.92), Inches(1.62), width=Inches(5.45), height=Inches(4.95))
    if right_img.exists():
        slide.shapes.add_picture(str(right_img), Inches(6.78), Inches(1.62), width=Inches(5.45), height=Inches(4.95))

    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_discussion_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, 'Discussion: 왜 어떤 경우는 이기고, 어떤 경우는 지는가')
    _add_logo(slide)
    if DIAGRAM_QUADRANT.exists():
        slide.shapes.add_picture(str(DIAGRAM_QUADRANT), Inches(0.88), Inches(1.42), width=Inches(7.2), height=Inches(5.35))
    body = slide.shapes.add_textbox(Inches(8.35), Inches(1.68), Inches(3.55), Inches(4.9))
    _set_text_frame(
        body.text_frame,
        [
            'small sparse: saved work 절대량이 작고 GPU fill도 낮다',
            'small dense: atom 수는 작아도 edge load가 충분하면 이길 수 있다',
            'large sparse: 크기만 크다고 자동으로 이기진 않지만 break-even을 넘기기 시작한다',
            'large dense: geometry reuse가 가장 크게 먹히는 구간이다',
        ],
        font_size=16,
        color=INK,
        bullet=True,
    )
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_status_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, '지금까지 구현한 것과 남은 것')
    _add_logo(slide)
    left = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.92), Inches(1.55), Inches(5.3), Inches(4.65))
    left.fill.solid()
    left.fill.fore_color.rgb = GREEN_SOFT
    left.line.color.rgb = STEEL
    ltitle = slide.shapes.add_textbox(Inches(1.1), Inches(1.72), Inches(4.6), Inches(0.3))
    p1 = ltitle.text_frame.paragraphs[0]
    p1.text = '완료'
    p1.font.size = Pt(22)
    p1.font.bold = True
    p1.font.color.rgb = INK
    lbody = slide.shapes.add_textbox(Inches(1.12), Inches(2.18), Inches(4.7), Inches(3.55))
    _set_text_frame(
        lbody.text_frame,
        [
            'pair metadata / topology cache 정리',
            'geometry + weight reuse 구현',
            'baseline / pair 실험 파이프라인',
            'size-density / profiling / GPU util 분석',
            '세미나용 자원 및 다이어그램 생성',
        ],
        font_size=18,
        color=INK,
        bullet=True,
    )

    right = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(6.8), Inches(1.55), Inches(5.3), Inches(4.65))
    right.fill.solid()
    right.fill.fore_color.rgb = ACCENT_SOFT
    right.line.color.rgb = STEEL
    rtitle = slide.shapes.add_textbox(Inches(6.98), Inches(1.72), Inches(4.6), Inches(0.3))
    p2 = rtitle.text_frame.paragraphs[0]
    p2.text = '남은 핵심 과제'
    p2.font.size = Pt(22)
    p2.font.bold = True
    p2.font.color.rgb = INK
    rbody = slide.shapes.add_textbox(Inches(7.0), Inches(2.18), Inches(4.7), Inches(3.55))
    _set_text_frame(
        rbody.text_frame,
        [
            'pair-major TP kernel',
            'FlashTP와의 진짜 backend co-design',
            'distributed pair schedule',
            '더 큰 public dataset 완전 확보',
            'additive profiler 정리와 논문화',
        ],
        font_size=18,
        color=INK,
        bullet=True,
    )
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_conclusion_slide(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, '결론')
    _add_logo(slide)

    card = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(1.05), Inches(1.8), Inches(11.0), Inches(3.9))
    card.fill.solid()
    card.fill.fore_color.rgb = STEEL_SOFT
    card.line.color.rgb = STEEL
    body = slide.shapes.add_textbox(Inches(1.32), Inches(2.12), Inches(10.45), Inches(3.1))
    _set_text_frame(
        body.text_frame,
        [
            '재사용 연산값을 pair-shared / parity-recoverable / direction-specific / node-global로 나누는 것이 출발점이다.',
            '현재 SevenNet 구현은 geometry/SH/weight reuse까지는 왔지만, pair-major TP는 아직 아니다.',
            '실험상 성패는 결국 edge load와 density의 조합으로 갈리고, 가장 강한 구간은 dense-large이다.',
            '다음 단계의 연구 포인트는 pair-major execution과 backend/runtime co-design이다.',
        ],
        font_size=20,
        color=INK,
        bullet=True,
    )
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def _add_appendix_download(prs, slide_no: int):
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    _add_background(slide)
    _title_box(slide, 'Appendix: 큰 데이터셋 확보 상태')
    _add_logo(slide)
    if PLOT_DOWNLOAD.exists():
        slide.shapes.add_picture(str(PLOT_DOWNLOAD), Inches(0.92), Inches(1.45), width=Inches(11.2), height=Inches(5.5))
    _add_footer(slide, 'Lab Seminar', slide_no)
    return slide


def build_deck():
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    _add_title_slide(prs, 1)
    _add_agenda_slide(prs, 2)
    _add_bullet_slide(
        prs,
        3,
        'GNN-IP란 무엇인가',
        [
            '원자 구조를 그래프로 보고 total energy를 예측한 뒤 force와 stress를 gradient로 얻는 모델',
            '핵심 입력은 원자 종류, 상대 위치, cutoff 안의 local neighborhood',
            '에너지 자체는 rotation-invariant scalar이고, force는 위치 미분으로 나온다',
            '따라서 intermediate representation은 equivariant하게 유지하는 것이 유리하다',
        ],
        note='이 슬라이드는 작업자 외 학생들을 위한 1-2장짜리 전체 설명 파트다.',
    )
    _add_pipeline_overview(prs, 4)
    _add_classification_slide(prs, 5)
    _add_current_vs_target_slide(prs, 6)
    _add_related_work_slide(prs, 7)
    _add_bullet_slide(
        prs,
        8,
        '실험 질문과 평가 프로토콜',
        [
            '질문 1: 어떤 그래프에서 pair execution이 실제로 이득인가?',
            '질문 2: 재사용되는 연산은 전체 시간에서 얼마나 큰 비중인가?',
            '질문 3: 현재 구현은 exact pair-major execution과 무엇이 다른가?',
            '환경: SevenNet / 7net-omni / single-GPU lab benchmark + public real dataset evaluation',
        ],
        note='속도, speedup, profiling, GPU utilization, size-density를 같이 봤다.',
    )
    _add_image_slide(
        prs,
        9,
        '데이터셋 size-density map',
        PLOT_MAP,
        bullets=[
            'x축: 대표 샘플 directed edges',
            'y축: 평균 directed neighbors',
            'threshold: large = edges >= 3000, dense = avg neighbors >= 40',
            '발표용으로 네 구역 representative를 명시해 비교했다',
        ],
        image_right=False,
    )
    _add_two_image_slide(
        prs,
        10,
        '네 구역 대표 비교',
        PLOT_LAT,
        PLOT_SPEED,
        'Latency comparison',
        'Speedup comparison',
    )
    _add_image_slide(
        prs,
        11,
        'Quadrant 결과의 핵심 해석',
        DIAGRAM_QUADRANT,
        bullets=[
            'small sparse: SPICE 2023, 0.743x',
            'small dense: phononDB, 1.105x',
            'large sparse: OMol25 validation, 1.126x',
            'large dense: MPtrj validation, 1.669x',
            '즉 atom 수 단독이 아니라 edge load와 density 조합이 핵심이다',
        ],
        image_right=True,
    )
    _add_image_slide(
        prs,
        12,
        'Baseline 상세 프로파일: extreme cases',
        PLOT_STAGE,
        bullets=[
            'small sparse baseline은 TP와 force_output 비중이 크다',
            'large dense baseline은 SH + TP + force_output이 동시에 크다',
            'weight_nn, gather, aggregation은 baseline에서도 상대적으로 작다',
            '그래서 현재 pair execution의 이득은 일부 구간에서만 크게 보인다',
        ],
        image_right=False,
    )
    _add_bullet_slide(
        prs,
        13,
        '랩 실험 결과 요약',
        [
            'steady-state 기준, dense-large 구간에서 speedup이 가장 크고 dense/small or sparse/small에서는 손해 또는 약한 이득',
            'GPU utilization과 stage profiling을 같이 보면, 작은 그래프는 saved work 절대량이 작고 under-fill 문제가 남는다',
            '반대로 dense-large에서는 geometry reuse가 실제 wall time 감소로 이어진다',
            '현재 구현은 pair-major TP가 아니므로 speedup ceiling이 존재한다',
        ],
        note='단일 GPU lab benchmark와 public real dataset benchmark를 함께 해석했다.',
    )
    _add_discussion_slide(prs, 14)
    _add_status_slide(prs, 15)
    _add_conclusion_slide(prs, 16)
    _add_appendix_download(prs, 17)

    prs.save(str(OUTFILE))


def write_script():
    lines = ['# Seminar_김민창 GNN-IP 최종 세미나 대본', '']
    for item in NOTES:
        lines.append(f"## {item['title']} ({item['time']})")
        lines.append(item['script'])
        lines.append('')
    SCRIPT_OUT.write_text('\n'.join(lines).strip() + '\n')


if __name__ == '__main__':
    build_deck()
    write_script()
