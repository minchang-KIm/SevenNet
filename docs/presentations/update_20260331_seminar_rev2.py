from __future__ import annotations

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.util import Pt


ROOT = Path(__file__).resolve().parent
PPT_PATH = ROOT / "20260331_Seminar_김민창_rev2.pptx"
SCRIPT_PATH = ROOT / "20260331_Seminar_김민창_rev2_대본.md"

INK = RGBColor(28, 34, 41)


SLIDE_SCRIPT = [
    (
        "Slide 1. Title",
        "1.0 min",
        "안녕하세요. 오늘 발표에서는 GNN-IP inference optimization reconstruction이라는 제목으로, "
        "지난주 계획서에서 제안했던 아이디어를 실제 SevenNet 구현과 실험 결과까지 연결해서 말씀드리겠습니다. "
        "핵심은 새로운 물리 모델을 제안하는 것이 아니라, equivariant interatomic potential이 실제로 어떤 계산 파이프라인으로 돌아가고 "
        "그 안에서 무엇을 정확히 재사용할 수 있는지 재구성하는 것입니다.",
    ),
    (
        "Slide 2. Table of Contents",
        "1.0 min",
        "발표는 이 순서로 진행하겠습니다. 먼저 GNN-IP 전체 맥락과 equivariant model pipeline을 짧게 보고, "
        "그다음 문제 정의와 관련 연구를 통해 왜 pair reconstruction이 필요한지 설명하겠습니다. "
        "이후 제안 방식과 실험 계획, 그리고 지금까지 확보한 단일 GPU 실험 결과와 discussion으로 마무리하겠습니다.",
    ),
    (
        "Slide 3. GNN-IP R&D Proposal",
        "2.0 min",
        "이 장은 전체 프로젝트 맥락입니다. 제가 맡은 주제는 Year 3의 GNN-IP kernel level optimization과 runtime optimization 사이에 걸쳐 있습니다. "
        "즉 tensor product 자체를 빠르게 만드는 문제도 중요하지만, 그 이전 단계의 non-TP geometric computation과 입력 표현을 다시 짜서 "
        "중복 계산을 없애는 문제가 동시에 중요합니다. 오늘 발표는 그중에서도 input representation restructuring 쪽에 초점을 둡니다.",
    ),
    (
        "Slide 4. Runtime / Distributed Roadmap",
        "1.5 min",
        "뒤의 연차 계획까지 보면 결국 목표는 단일 커널 최적화에 그치지 않습니다. "
        "적응적 skipping 정책, runtime optimization, 그리고 multi-GPU, multi-node inference까지 연결되어 있습니다. "
        "그래서 이번 pair reconstruction도 단순한 micro-optimization이 아니라, 이후 distributed runtime으로 확장 가능한 "
        "중간 표현을 정의하는 문제로 보는 것이 맞습니다.",
    ),
    (
        "Slide 5. Equivariant Model Pipeline",
        "2.5 min",
        "NequIP 기준 전체 파이프라인을 먼저 짚고 가겠습니다. 원자 좌표와 원자종에서 neighbor list를 만들고, "
        "relative vector로부터 distance, radial basis, cutoff, spherical harmonics를 계산해서 edge feature를 만듭니다. "
        "그 후 각 interaction block에서 weight neural network, tensor product, aggregation이 반복되고, "
        "마지막에 energy를 읽고 force와 stress를 gradient로 얻습니다. "
        "오늘 이야기할 재사용 가능성은 거의 전부 이 파이프라인의 edge geometry와 message passing 사이에서 결정됩니다.",
    ),
    (
        "Slide 6. Problem",
        "2.5 min",
        "문제 정의는 간단합니다. directed edge 기반 구현에서는 i에서 j로 가는 edge와 j에서 i로 가는 edge가 별도로 처리됩니다. "
        "그런데 이 둘은 거리, radial basis, cutoff처럼 완전히 같은 값과, spherical harmonics처럼 sign 규칙으로 복원 가능한 값을 많이 공유합니다. "
        "반면 message와 aggregation은 source node가 달라서 그대로 공유할 수 없습니다. "
        "즉 reusable과 non-reusable을 정확히 분류해야 원래 formulation을 깨지 않고 어느 부분까지 줄일 수 있는지가 보입니다.",
    ),
    (
        "Slide 7. Challenge 1",
        "2.0 min",
        "첫 번째 challenge는 structural inefficiency입니다. 기존 edge-wise flow는 i to j와 j to i를 완전히 독립으로 취급해서 "
        "geometric computation을 두 번 합니다. 그래서 pair-wise reconstructuring이 필요합니다. "
        "다만 그냥 pair로 묶는다고 끝나는 것이 아니라, sign flip으로 정확히 재사용할 수 있는 값과 "
        "방향별로 다시 계산해야 하는 값을 구분해야 합니다.",
    ),
    (
        "Slide 8. Challenge 2",
        "2.0 min",
        "두 번째 challenge는 backend integration입니다. FlashTP 같은 시스템은 tensor product를 매우 잘 가속하지만, "
        "현재 인터페이스는 여전히 directed edge 중심입니다. 그래서 pair-based input을 설계하더라도 실제로 이득을 보려면 "
        "backend API와 kernel layout까지 같이 고려해야 합니다. "
        "이 지점 때문에 단순한 파이썬 레벨 refactoring과 pair-major execution은 난이도가 크게 다릅니다.",
    ),
    (
        "Slide 9. Related Work",
        "2.5 min",
        "관련 연구는 세 갈래로 보면 충분합니다. FlashTP와 SevenNet은 tensor product와 시스템 측면을 강하게 최적화하지만, "
        "bidirectional geometry redundancy는 거의 건드리지 않습니다. "
        "반대로 NEMP는 node space로 옮겨가며 효율을 얻지만, 원래 E3-equivariant message passing formulation 자체를 바꿉니다. "
        "제가 찾은 범위에서는, 원래 formulation과 파라미터를 유지한 채 bidirectional geometric redundancy를 정면으로 줄이는 접근은 많지 않았습니다.",
    ),
    (
        "Slide 10. Proposed Method Motivation",
        "2.0 min",
        "이 장은 왜 sign-flip reuse가 수학적으로 가능한지를 보여주는 배경입니다. "
        "E3-equivariant 모델은 rotation, translation, reflection에 대해 정해진 방식으로 변해야 합니다. "
        "translation은 relative coordinate로 처리되고, reflection은 spherical harmonics의 parity, 즉 degree l에 따른 부호 변화로 처리됩니다. "
        "그래서 역방향 edge의 SH를 새로 계산하지 않고 parity sign으로 복원할 수 있다는 아이디어가 나옵니다.",
    ),
    (
        "Slide 11. Reason Why choose this",
        "2.5 min",
        "이 페이지를 고른 이유는 두 가지입니다. 첫째, SH는 이론적으로 재사용이 가장 명확한 항목입니다. "
        "둘째, 실제 profiling에서도 large dense case에서는 무시할 수 없는 비용입니다. "
        "저희가 baseline을 intrusive profiling한 결과, MPtrj와 MD22 nanotube 같은 큰 dense 그래프에서는 "
        "spherical harmonics가 전체 model time의 약 18퍼센트에서 26퍼센트를 차지했습니다. "
        "즉 모든 데이터셋에서 지배적이라고 말할 수는 없지만, 큰 edge load 조건에서는 분명히 substantial한 비중입니다. "
        "그래서 이 값은 compute once, reverse reuse의 첫 번째 타깃으로 적합합니다.",
    ),
    (
        "Slide 12. Proposed Method",
        "3.0 min",
        "현재 제가 구현한 방식은 pair metadata를 먼저 만들고, 하나의 undirected pair를 기준으로 geometry를 한 번만 계산하는 것입니다. "
        "distance, radial basis, cutoff, spherical harmonics를 pair 기준으로 계산하고, reverse edge의 SH는 parity sign으로 복원합니다. "
        "또 weight neural network도 pair 기준으로 한 번 계산해 재사용합니다. "
        "하지만 여기서 중요한 한계가 있습니다. 최종 tensor product와 aggregation은 아직 directed edge 기준으로 남아 있습니다. "
        "즉 지금 단계는 pair-aware reuse이지, 완전한 pair-major TP kernel은 아닙니다.",
    ),
    (
        "Slide 13. Experiment Plan",
        "3.5 min",
        "실험은 현재 두 층으로 진행하고 있습니다. 첫째는 SevenNet e3nn baseline과 SevenNet plus pair execution full을 직접 비교하는 단일 GPU benchmark입니다. "
        "둘째는 FlashTP reference를 함께 놓고, 현재 구현이 geometry reuse 수준인지 아니면 backend까지 이득을 주는지 확인하는 비교입니다. "
        "워크로드는 small sparse, small dense, large sparse, large dense 네 구역을 모두 포함하도록 잡았습니다. "
        "대표적으로 SPICE와 ANI-1x, phononDB, OMol25 validation, MPtrj, MD22 nanotube 등을 봤습니다. "
        "지표는 정확도 차이뿐 아니라 cold와 steady latency, stage profiling, GPU utilization, size-density map까지 함께 측정했습니다. "
        "현재까지는 단일 GPU 결과가 정리되어 있고, multi-GPU scaling은 다음 단계 과제로 남아 있습니다.",
    ),
    (
        "Slide 14. Discussion",
        "2.0 min",
        "정리하면, 현재 결과는 과장 없이 이렇게 보는 것이 맞습니다. "
        "현대 SH 구현은 느린 삼각함수 기반이 아니라 Cartesian polynomial 기반이지만, edge 수가 크면 여전히 무시할 수 없는 비용입니다. "
        "현재 구현은 geometry와 SH, weight reuse에서 이득을 얻지만, TP와 aggregation은 그대로 남아 있기 때문에 large dense graph에서만 확실히 이깁니다. "
        "따라서 다음 단계는 pair-major TP kernel과 backend co-design이며, 그래야 FlashTP 대비로도 일관된 end-to-end acceleration을 주장할 수 있습니다.",
    ),
]


def _set_text_lines(shape, lines: list[str]) -> None:
    tf = shape.text_frame
    paragraphs = list(tf.paragraphs)
    while len(paragraphs) < len(lines):
        paragraphs.append(tf.add_paragraph())
    for i, line in enumerate(lines):
        paragraphs[i].text = line
    for j in range(len(lines), len(paragraphs)):
        paragraphs[j].text = ""


def _style_textbox(textbox) -> None:
    tf = textbox.text_frame
    for i, p in enumerate(tf.paragraphs):
        for run in p.runs:
            run.font.size = Pt(22 if i % 2 == 0 else 19)
            run.font.color.rgb = INK
            run.font.bold = (i % 2 == 0)


def update_ppt() -> None:
    prs = Presentation(str(PPT_PATH))

    slide11 = prs.slides[10]
    _set_text_lines(
        slide11.shapes[1],
        [
            "1. The Bottleneck: Spherical Harmonics",
            "Even with modern Cartesian-polynomial implementations, SH remains a visible geometric cost when the graph is large.",
            "Profiling Insight: In our lab baseline profiling, SH takes roughly 18-26% of model time on large dense cases such as MPtrj and MD22 nanotube.",
            "2. Mathematical Foundation: Parity Property",
            "The inputs for bidirectional edges (i to j and j to i) contain identical geometric information, differing only by an equivariant transformation.",
            "3. Proposed Solution: Sign Flip Reuse",
            "Compute Once: Calculate spherical harmonics only for one canonical direction within each pair.",
            "Instant Reuse: Recover the reverse edge by parity sign, instead of recomputing the same geometry.",
        ],
    )

    slide12 = prs.slides[11]
    if len(slide12.shapes) < 3 or not getattr(slide12.shapes[2], "has_text_frame", False):
        textbox = slide12.shapes.add_textbox(205740, 927088, 11757660, 5146052)
    else:
        textbox = slide12.shapes[2]
    _set_text_lines(
        textbox,
        [
            "1. Pair Metadata Construction",
            "Group (i to j) and (j to i) into one undirected pair and keep forward / reverse indices explicitly.",
            "2. Pair-shared Geometry",
            "Compute distance, radial basis, cutoff, and spherical harmonics once per pair instead of once per directed edge.",
            "3. Exact Reverse Recovery",
            "Recover reverse spherical harmonics with parity sign, not with a second SH evaluation.",
            "4. Current Implementation Boundary",
            "weight_nn is also reused at pair level, but tensor product and aggregation are still executed on directed edges.",
        ],
    )
    _style_textbox(textbox)

    slide13 = prs.slides[12]
    _set_text_lines(
        slide13.shapes[1],
        [
            "Baselines",
            "SevenNet e3nn baseline / SevenNet + pair execution(full) / FlashTP reference",
            "Workloads",
            "Small sparse: SPICE 2023, ANI-1x",
            "Small dense: phononDB, selected periodic crystals",
            "Large sparse: OMol25 validation, OC20 validation",
            "Large dense: MPtrj, MD22 nanotube, OMat24",
            "Evaluation Metrics",
            "Energy-force exactness",
            "Cold / steady latency + stage profiling + GPU utilization",
            "Size-density analysis and quadrant comparison",
            "Scope: single-GPU lab benchmark first, multi-GPU scaling next",
        ],
    )

    slide14 = prs.slides[13]
    _set_text_lines(
        slide14.shapes[1],
        [
            "Modern SH Implementation",
            "Uses fast Cartesian polynomials, so the win is not from replacing slow trigonometric code.",
            "Current Bottleneck",
            "The real gain comes from removing duplicated geometry, SH, and pair-weight work on reverse edges.",
            "Current Limitation",
            "TP and aggregation are still directed-edge operations, so the benefit is strongest only on large dense graphs.",
            "Next Step",
            "Pair-major TP plus backend co-design is required before claiming consistent end-to-end acceleration over FlashTP.",
        ],
    )

    prs.save(str(PPT_PATH))


def write_script() -> None:
    lines = ["# 20260331 Seminar 김민창 rev2 대본", ""]
    for title, time, script in SLIDE_SCRIPT:
        lines.append(f"## {title} ({time})")
        lines.append(script)
        lines.append("")
    SCRIPT_PATH.write_text("\n".join(lines), encoding="utf-8")


if __name__ == "__main__":
    update_ppt()
    write_script()
